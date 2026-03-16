"""
Phase 1: PPTX → raw JSON extraction (no API calls).

Walks a POWERPOINTS folder tree, extracts text from each game's PPTX,
optionally looks up the game in market_names.xlsx, and writes one JSON
per game to the output directory.
"""

import json
import re
import logging
from pathlib import Path

from pptx import Presentation
from rapidfuzz import fuzz
from tqdm import tqdm

logger = logging.getLogger(__name__)

# Numeric-prefix pattern for game folders: 0001_GameName
GAME_FOLDER_RE = re.compile(r'^(\d{3,4})_(.+)$')

# Market suffix pattern for stripping tablename → base_key
SUFFIX_PATTERN = re.compile(r'(NoIp|Es|Pt|It|Co|Nl|Ca|Se)$', re.IGNORECASE)

# Noise patterns to skip during text extraction
SKIP_PATTERNS = [
    re.compile(r'^\*'),                          # disclaimer lines
    re.compile(r'^\d+$'),                        # standalone numbers (layout artefacts)
    re.compile(r'Gráficos cogidos', re.IGNORECASE),
    re.compile(r'Gr.ficos cogidos', re.IGNORECASE),   # encoding variants
    re.compile(r'Gráficos no finales', re.IGNORECASE),
    re.compile(r'Gr.ficos no finales', re.IGNORECASE),
    re.compile(r'Localización \(Audios', re.IGNORECASE),
    re.compile(r'Localizaci.n \(Audios', re.IGNORECASE),
    re.compile(r'Internacional MGA', re.IGNORECASE),
]

# Subdirectories to skip when looking for PPTX files
SKIP_DIRS = {'old', 'aprobaciones'}


def _should_skip_text(text: str) -> bool:
    """Return True if this text line is noise that should be filtered out."""
    return any(p.search(text) for p in SKIP_PATTERNS)


def _find_game_folders(pptx_root: Path) -> list[dict]:
    """
    Walk pptx_root and find all game folders (matching NNNN_GameName pattern).
    Returns list of dicts with folder_path, folder_category, folder_game_name, game_id.

    Handles varying depth: MegaWays/NNNN_Game or Slots5/Slots5/NNNN_Game.
    """
    games = []
    for path in sorted(pptx_root.rglob('*')):
        if not path.is_dir():
            continue
        match = GAME_FOLDER_RE.match(path.name)
        if not match:
            continue
        # Skip if this is inside an old/aprobaciones subfolder
        parts_lower = [p.lower() for p in path.parts]
        if any(skip in parts_lower for skip in SKIP_DIRS):
            continue

        game_id = match.group(1)
        game_name = match.group(2)

        # Derive category from the first directory level under pptx_root
        try:
            rel = path.relative_to(pptx_root)
            category = rel.parts[0] if rel.parts else 'UNKNOWN'
        except ValueError:
            category = 'UNKNOWN'

        games.append({
            'folder_path': path,
            'folder_category': category,
            'folder_game_name': game_name,
            'game_id': game_id,
        })
    return games


def _select_pptx(game_folder: Path) -> Path | None:
    """
    Select the best PPTX file from a game folder.
    Priority: filename containing "Comercial" (case-insensitive) > first .pptx alphabetically.
    Skips files inside old/aprobaciones subdirs.
    """
    pptx_files = []
    for f in sorted(game_folder.rglob('*.pptx')):
        # Skip files in old/aprobaciones subdirs
        rel = f.relative_to(game_folder)
        parts_lower = [p.lower() for p in rel.parts[:-1]]  # exclude filename
        if any(skip in parts_lower for skip in SKIP_DIRS):
            continue
        pptx_files.append(f)

    if not pptx_files:
        return None

    # Prefer Comercial files
    for f in pptx_files:
        if 'comercial' in f.name.lower():
            return f

    # Fallback: first alphabetically
    return pptx_files[0]


def _extract_slides(pptx_path: Path) -> list[dict]:
    """
    Extract text from all slides in a PPTX file.
    Returns list of slide dicts with slide_num, is_category_slide, is_sales_slide, texts.
    """
    slides_data = []
    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        logger.warning(f"Failed to open {pptx_path}: {e}")
        return slides_data

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and not _should_skip_text(text):
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text and not _should_skip_text(text):
                            texts.append(text)

        all_text_joined = ' '.join(texts).upper()
        is_category = bool(
            'CATEGORIAS' in all_text_joined
            or 'TEMATICAS' in all_text_joined
            or 'TEMÁTICAS' in all_text_joined
        )
        is_sales = bool(
            ('ARGUMENTOS' in all_text_joined and 'VENTA' in all_text_joined)
            or 'CARACTERISTICAS' in all_text_joined
            or 'CARACTERÍSTICAS' in all_text_joined
        )

        # Filter out DESARROLLO section content on category slides
        filtered_texts = []
        in_desarrollo = False
        for text in texts:
            if text.upper().strip() == 'DESARROLLO':
                in_desarrollo = True
                continue
            # End DESARROLLO section at the next recognisable header or theme keyword
            if in_desarrollo:
                # DESARROLLO content is typically a single block — skip until we hit
                # a known theme/feature keyword or short tag-like line
                if len(text) > 100 or text.lower().startswith('inspirad'):
                    continue
                in_desarrollo = False
            filtered_texts.append(text)

        slides_data.append({
            'slide_num': i,
            'is_category_slide': is_category,
            'is_sales_slide': is_sales,
            'texts': filtered_texts,
        })

    return slides_data


def _load_market_db(market_names_path: Path) -> list[dict] | None:
    """Load market_names.xlsx if it exists. Returns list of row dicts or None."""
    if not market_names_path.exists():
        logger.info(f"Market names file not found at {market_names_path} — skipping market lookup")
        return None

    import openpyxl
    wb = openpyxl.load_workbook(str(market_names_path), read_only=True)
    ws = wb.active
    headers = [cell.value for cell in next(ws.iter_rows(min_row=1, max_row=1))]
    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        row_dict = dict(zip(headers, row))
        rows.append(row_dict)
    wb.close()
    return rows


def _lookup_market(game_name: str, market_rows: list[dict]) -> dict:
    """
    Match a game folder name against market_names.xlsx rows.
    Returns match info dict.
    """
    # Clean game_name: replace underscores with spaces, strip
    clean_name = game_name.replace('_', ' ').strip()

    # 1. Exact match on 'name' column
    for row in market_rows:
        name = str(row.get('name', ''))
        if name.lower() == clean_name.lower():
            return _build_match_result(row, market_rows, 1.0, 'exact_name')

    # 2. Exact match on 'COMMERCIAL NAME'
    for row in market_rows:
        cname = str(row.get('COMMERCIAL NAME', ''))
        if cname.lower() == clean_name.lower():
            return _build_match_result(row, market_rows, 1.0, 'exact_commercial')

    # 3. Fuzzy match on COMMERCIAL NAME
    best_score = 0
    best_row = None
    for row in market_rows:
        cname = str(row.get('COMMERCIAL NAME', ''))
        score = fuzz.token_sort_ratio(clean_name.lower(), cname.lower())
        if score > best_score:
            best_score = score
            best_row = row

    if best_score >= 75 and best_row:
        return _build_match_result(best_row, market_rows, best_score / 100.0, 'fuzzy_commercial')

    # 4. Fuzzy match on name column
    best_score = 0
    best_row = None
    for row in market_rows:
        name = str(row.get('name', ''))
        score = fuzz.token_sort_ratio(clean_name.lower(), name.lower())
        if score > best_score:
            best_score = score
            best_row = row

    if best_score >= 75 and best_row:
        return _build_match_result(best_row, market_rows, best_score / 100.0, 'fuzzy_name')

    return {
        'matched_name': None,
        'matched_commercial_name': None,
        'matched_tablename': None,
        'base_key': None,
        'all_markets': [],
        'match_confidence': 0.0,
        'match_method': 'none',
    }


def _build_match_result(row: dict, all_rows: list[dict], confidence: float, method: str) -> dict:
    """Build a market match result from a matched row."""
    tablename = str(row.get('tablename', ''))
    base_key = SUFFIX_PATTERN.sub('', tablename) if tablename else None

    # Find all markets for this game family
    all_markets = []
    if base_key:
        for r in all_rows:
            tn = str(r.get('tablename', ''))
            if SUFFIX_PATTERN.sub('', tn) == base_key:
                market = r.get('MARKET', '')
                if market and market not in all_markets:
                    all_markets.append(market)

    return {
        'matched_name': row.get('name'),
        'matched_commercial_name': row.get('COMMERCIAL NAME'),
        'matched_tablename': tablename,
        'base_key': base_key,
        'all_markets': all_markets,
        'match_confidence': round(confidence, 3),
        'match_method': method,
    }


def extract_all(
    pptx_root: Path,
    market_names_path: Path,
    output_dir: Path,
) -> dict:
    """
    Main extraction function. Walks pptx_root, extracts text from PPTXs,
    optionally matches against market_names.xlsx, writes one JSON per game.

    Returns a summary dict with stats.
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    # Find all game folders
    game_folders = _find_game_folders(pptx_root)
    logger.info(f"Found {len(game_folders)} game folders")

    # Load market DB if available
    market_rows = _load_market_db(market_names_path)

    stats = {
        'total_folders': len(game_folders),
        'pptx_found': 0,
        'pptx_not_found': 0,
        'market_exact': 0,
        'market_fuzzy': 0,
        'market_none': 0,
        'market_skipped': 0,
        'json_written': 0,
        'errors': 0,
    }

    for game_info in tqdm(game_folders, desc="Extracting games"):
        folder_path = game_info['folder_path']
        game_name = game_info['folder_game_name']
        category = game_info['folder_category']
        game_id = game_info['game_id']

        # Select PPTX
        pptx_path = _select_pptx(folder_path)
        pptx_found = pptx_path is not None

        if pptx_found:
            stats['pptx_found'] += 1
        else:
            stats['pptx_not_found'] += 1

        # Extract slides
        slides_data = []
        if pptx_found:
            try:
                slides_data = _extract_slides(pptx_path)
            except Exception as e:
                logger.error(f"Error extracting {pptx_path}: {e}")
                stats['errors'] += 1

        # Market lookup
        if market_rows is not None:
            market_lookup = _lookup_market(game_name, market_rows)
            method = market_lookup['match_method']
            if method.startswith('exact'):
                stats['market_exact'] += 1
            elif method.startswith('fuzzy'):
                stats['market_fuzzy'] += 1
            else:
                stats['market_none'] += 1
        else:
            stats['market_skipped'] += 1
            market_lookup = {
                'matched_name': None,
                'matched_commercial_name': None,
                'matched_tablename': None,
                'base_key': None,
                'all_markets': [],
                'match_confidence': 0.0,
                'match_method': 'skipped',
            }

        # Use base_key from market lookup, or derive from folder name
        base_key = market_lookup.get('base_key')
        if not base_key:
            # Derive a key from game_id + cleaned game name
            base_key = f"{category}_{game_id}_{game_name}".replace(' ', '_')

        # Build output record
        record = {
            'base_key': base_key,
            'game_id': game_id,
            'folder_category': category,
            'folder_game_name': game_name,
            'pptx_filename': pptx_path.name if pptx_path else None,
            'pptx_found': pptx_found,
            'market_lookup': market_lookup,
            'slides': slides_data,
        }

        # Write JSON
        safe_name = re.sub(r'[^\w\-]', '_', base_key)
        output_path = output_dir / f"{safe_name}.json"
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(record, f, indent=2, ensure_ascii=False)
            stats['json_written'] += 1
        except Exception as e:
            logger.error(f"Error writing {output_path}: {e}")
            stats['errors'] += 1

    return stats
